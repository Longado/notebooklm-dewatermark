#!/usr/bin/env python3
"""NotebookLM PPT Watermark Remover

Remove the NotebookLM watermark (logo + text) from the bottom-right corner
of slides in exported PPTX files.

Usage:
    python3 notebooklm_dewatermark.py input.pptx
    python3 notebooklm_dewatermark.py input.pptx -o output.pptx
    python3 notebooklm_dewatermark.py *.pptx
    python3 notebooklm_dewatermark.py input.pptx --wm-width 200 --wm-height 40
"""

import argparse
import io
import sys
from pathlib import Path

from PIL import Image, ImageFilter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Default watermark region parameters (offset from bottom-right corner)
DEFAULT_WM_WIDTH = 175
DEFAULT_WM_HEIGHT = 30
DEFAULT_MARGIN_RIGHT = 3
DEFAULT_MARGIN_BOTTOM = 3
SAMPLE_PADDING = 5

# Minimum coverage ratio to consider a picture as a full-slide background
MIN_COVERAGE = 0.85

# Feather width in pixels for edge blending
FEATHER_PX = 6

# XML namespaces
NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
NS_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def has_watermark(img, wm_left, wm_top, wm_right, wm_bottom):
    """Detect whether a watermark likely exists in the given region.

    Compares the pixel variance inside the watermark region against the
    surrounding background. If the watermark region has notably different
    content than the background strip above it, a watermark is likely present.
    """
    w, h = img.size

    # Sample a background reference strip above the watermark
    ref_top = max(0, wm_top - (wm_bottom - wm_top) - SAMPLE_PADDING)
    ref_bottom = wm_top
    ref_strip = img.crop((wm_left, ref_top, wm_right, ref_bottom))
    wm_strip = img.crop((wm_left, wm_top, wm_right, wm_bottom))

    ref_pixels = list(ref_strip.getdata())
    wm_pixels = list(wm_strip.getdata())

    if not ref_pixels or not wm_pixels:
        return False

    def avg_color(pixels):
        channels = len(pixels[0]) if isinstance(pixels[0], tuple) else 1
        if channels == 1:
            return (sum(pixels) / len(pixels),)
        sums = [0] * channels
        for px in pixels:
            for c in range(channels):
                sums[c] += px[c]
        return tuple(s / len(pixels) for s in sums)

    ref_avg = avg_color(ref_pixels)
    wm_avg = avg_color(wm_pixels)

    # Compare average color difference across RGB channels
    diff = sum(abs(r - w) for r, w in zip(ref_avg[:3], wm_avg[:3])) / 3
    return diff > 3.0


def remove_watermark(img, wm_width, wm_height, margin_right, margin_bottom):
    """Remove the NotebookLM watermark from the bottom-right corner.

    Strategy: copy a strip of pixels from directly above the watermark area
    and paste it over the watermark, with feathered edges for seamless blending.
    """
    result = img.copy()
    w, h = result.size

    wm_left = w - wm_width - margin_right
    wm_top = h - wm_height - margin_bottom
    wm_right = w - margin_right
    wm_bottom = h - margin_bottom

    # Check if watermark exists
    if not has_watermark(img, wm_left, wm_top, wm_right, wm_bottom):
        return None

    sample_top = max(0, wm_top - wm_height - SAMPLE_PADDING)
    sample_bottom = wm_top

    sample_strip = result.crop((wm_left, sample_top, wm_right, sample_bottom))

    target_w = wm_right - wm_left
    target_h = wm_bottom - wm_top
    sample_strip = sample_strip.resize((target_w, target_h), Image.LANCZOS)

    # Create a feathered alpha mask for smooth blending
    mask = Image.new("L", (target_w, target_h), 255)
    for y in range(min(FEATHER_PX, target_h)):
        alpha = int(255 * y / FEATHER_PX)
        for x in range(target_w):
            mask.putpixel((x, y), alpha)

    # Also feather left edge
    for x in range(min(FEATHER_PX, target_w)):
        alpha_x = int(255 * x / FEATHER_PX)
        for y in range(target_h):
            current = mask.getpixel((x, y))
            mask.putpixel((x, y), min(current, alpha_x))

    # Ensure sample_strip matches result mode for compositing
    if result.mode == "RGBA" and sample_strip.mode != "RGBA":
        sample_strip = sample_strip.convert("RGBA")
    elif result.mode == "RGB" and sample_strip.mode != "RGB":
        sample_strip = sample_strip.convert("RGB")

    result.paste(sample_strip, (wm_left, wm_top), mask)

    return result


def is_fullpage_image(shape, slide_width, slide_height):
    """Check if a picture shape covers most of the slide area."""
    w_ratio = shape.width / slide_width
    h_ratio = shape.height / slide_height
    return w_ratio >= MIN_COVERAGE and h_ratio >= MIN_COVERAGE


def process_pptx(input_path, output_path, wm_width, wm_height,
                 margin_right, margin_bottom):
    """Process a single PPTX file. Returns (processed, skipped) counts."""
    prs = Presentation(str(input_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    processed = 0
    skipped = 0
    processed_part_ids = set()

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            # P0: Only process full-page background images
            if not is_fullpage_image(shape, slide_w, slide_h):
                skipped += 1
                continue

            # P0: Error handling for image access
            try:
                img_blob = shape.image.blob
                content_type = shape.image.content_type
            except Exception as e:
                print(f"  Warning: slide {slide_num}: cannot read image: {e}",
                      file=sys.stderr)
                skipped += 1
                continue

            # P0: Error handling for XML traversal
            try:
                pic = shape._element
                blip = pic.find(f".//{NS_A}blip")
                if blip is None:
                    skipped += 1
                    continue
                r_id = blip.get(f"{NS_R}embed")
                if r_id is None:
                    skipped += 1
                    continue
                image_part = slide.part.rels[r_id].target_part
            except Exception as e:
                print(f"  Warning: slide {slide_num}: cannot resolve image part: {e}",
                      file=sys.stderr)
                skipped += 1
                continue

            # P1: Skip already-processed image parts (shared across slides)
            part_id = id(image_part)
            if part_id in processed_part_ids:
                processed += 1
                continue
            processed_part_ids.add(part_id)

            try:
                img = Image.open(io.BytesIO(img_blob))
            except Exception as e:
                print(f"  Warning: slide {slide_num}: cannot decode image: {e}",
                      file=sys.stderr)
                skipped += 1
                continue

            # P1: Remove watermark with detection
            cleaned = remove_watermark(img, wm_width, wm_height,
                                       margin_right, margin_bottom)
            if cleaned is None:
                skipped += 1
                continue

            # Save with correct format, preserving alpha
            buf = io.BytesIO()
            if "png" in content_type:
                fmt = "PNG"
            elif img.mode == "RGBA":
                fmt = "PNG"
            else:
                fmt = "JPEG"
            cleaned.save(buf, format=fmt)
            buf.seek(0)

            image_part._blob = buf.read()
            processed += 1

    prs.save(str(output_path))
    return processed, skipped


def main():
    parser = argparse.ArgumentParser(
        description="Remove NotebookLM watermark from PPTX slides"
    )
    parser.add_argument(
        "input", nargs="+", type=Path, help="Input PPTX file(s)"
    )
    parser.add_argument(
        "-o", "--output", type=Path, default=None,
        help="Output file path (only for single file input)"
    )
    parser.add_argument(
        "--wm-width", type=int, default=DEFAULT_WM_WIDTH,
        help=f"Watermark width in pixels (default: {DEFAULT_WM_WIDTH})"
    )
    parser.add_argument(
        "--wm-height", type=int, default=DEFAULT_WM_HEIGHT,
        help=f"Watermark height in pixels (default: {DEFAULT_WM_HEIGHT})"
    )
    parser.add_argument(
        "--margin-right", type=int, default=DEFAULT_MARGIN_RIGHT,
        help=f"Right margin in pixels (default: {DEFAULT_MARGIN_RIGHT})"
    )
    parser.add_argument(
        "--margin-bottom", type=int, default=DEFAULT_MARGIN_BOTTOM,
        help=f"Bottom margin in pixels (default: {DEFAULT_MARGIN_BOTTOM})"
    )
    args = parser.parse_args()

    for input_path in args.input:
        if not input_path.exists():
            print(f"Error: file not found: {input_path}", file=sys.stderr)
            continue

        if input_path.suffix.lower() != ".pptx":
            print(f"Skipping non-PPTX file: {input_path}", file=sys.stderr)
            continue

        if args.output and len(args.input) == 1:
            output_path = args.output
        else:
            output_path = input_path.with_stem(input_path.stem + "_clean")

        processed, skipped = process_pptx(
            input_path, output_path,
            args.wm_width, args.wm_height,
            args.margin_right, args.margin_bottom,
        )
        status = f"Done: {input_path.name} -> {output_path.name}"
        status += f" ({processed} cleaned"
        if skipped:
            status += f", {skipped} skipped"
        status += ")"
        print(status)


if __name__ == "__main__":
    main()
